import win32com.client as win32
import time
import qi

from pptx import Presentation

def main(session):
    tts = session.service("ALTextToSpeech")
    tts.setVoice("Alyona22Enhanced")
    filePath = "C:\\Users\\tsi_nao\\Desktop\\TESTTHIS.pptx"
    ppoint = win32.gencache.EnsureDispatch('Powerpoint.Application')
    ppoint.Visible = True
    presentation = ppoint.Presentations.Open(filePath)
    presentation.SlideShowSettings.ShowWithAnimation = False
    slideShow = presentation.SlideShowSettings.Run()


    ppt=Presentation(filePath)
    notes = []



    for page, slide in enumerate(ppt.slides):
        textNote = slide.notes_slide.notes_text_frame.text
        notes.append(textNote)
        tts.say(str(notes[page].encode('utf-8')))

        slideShow.View.Next()
        time.sleep(1)

    presentation.Close()
    ppoint.Quit()



if __name__ == "__main__":
    session = qi.Session()
    session.connect("tcp://" + "192.168.252.226" + ":" + "9559")
    main(session)